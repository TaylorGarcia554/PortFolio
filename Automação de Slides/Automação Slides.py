import os
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog, messagebox

def change_background(pptx_path, new_background_image_path, output_path):
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        background = slide.shapes.add_picture(new_background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
    
    prs.save(output_path)

def process_folder(folder_path, new_background_image_path, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx") and not filename.startswith("~$"):
            pptx_path = os.path.join(folder_path, filename)
            output_path = os.path.join(output_folder, filename)
            try:
                change_background(pptx_path, new_background_image_path, output_path)
                print(f"Processed: {pptx_path}")
            except Exception as e:
                print(f"Erro ao processar {pptx_path}: {e}")

def select_image():
    image_selected = filedialog.askopenfilename(filetypes=[("PNG files", "*.png"), ("All files", "*.*")])
    if image_selected:
        image_path_entry.delete(0, tk.END)
        image_path_entry.insert(0, image_selected)

def select_output_folder():
    folder_selected = filedialog.askdirectory()
    if folder_selected:
        output_folder_entry.delete(0, tk.END)
        output_folder_entry.insert(0, folder_selected)

def start_processing():
    folder_path = "L:/SLIDES_2024"
    new_background_image_path = image_path_entry.get()
    output_folder = output_folder_entry.get()
    
    if not folder_path or not new_background_image_path or not output_folder:
        messagebox.showerror("Error", "Todos os campos devem ser preenchidos")
        return
    
    process_folder(folder_path, new_background_image_path, output_folder)
    
    image_path_entry.delete(0, tk.END)
    output_folder_entry.delete(0, tk.END)

# Configurar a interface Tkinter
root = tk.Tk()
root.title("Processador de Slides")

tk.Label(root, text="Imagem de fundo:").grid(row=1, column=0, padx=10, pady=5)
image_path_entry = tk.Entry(root, width=50)
image_path_entry.grid(row=1, column=1, padx=10, pady=5)
tk.Button(root, text="Selecionar", command=select_image).grid(row=1, column=2, padx=10, pady=5)

tk.Label(root, text="Pasta de saída:").grid(row=2, column=0, padx=10, pady=5)
output_folder_entry = tk.Entry(root, width=50)
output_folder_entry.grid(row=2, column=1, padx=10, pady=5)
tk.Button(root, text="Selecionar", command=select_output_folder).grid(row=2, column=2, padx=10, pady=5)

tk.Button(root, text="Iniciar Processamento", command=start_processing).grid(row=3, column=0, columnspan=3, pady=20)



root.mainloop()